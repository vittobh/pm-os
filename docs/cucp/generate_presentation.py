import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation(output_path):
    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Colors
    COLOR_BG = RGBColor(248, 250, 252)       # #F8FAFC
    COLOR_DARK = RGBColor(23, 43, 77)        # #172B4D
    COLOR_MUTED = RGBColor(94, 108, 132)     # #5E6C84
    COLOR_BLUE = RGBColor(0, 82, 204)        # #0052CC
    COLOR_BLUE_LIGHT = RGBColor(222, 235, 255) # #DEEBFF
    COLOR_RED = RGBColor(255, 86, 48)        # #FF5630
    COLOR_GREEN = RGBColor(54, 179, 126)     # #36B37E
    COLOR_GOLD = RGBColor(255, 171, 0)       # #FFAB00
    COLOR_GOLD_LIGHT = RGBColor(255, 240, 181) # #FFF0B5
    COLOR_WHITE = RGBColor(255, 255, 255)

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_header(slide, title_text, subtitle_text=""):
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.7))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Segoe UI'
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = COLOR_DARK
        
        # Subtitle
        if subtitle_text:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.4))
            tf_sub = sub_box.text_frame
            tf_sub.word_wrap = True
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle_text
            p_sub.font.name = 'Segoe UI'
            p_sub.font.size = Pt(13)
            p_sub.font.color.rgb = COLOR_MUTED

    slide_layout = prs.slide_layouts[6] # Blank layout

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide1)

    # Label
    lbl_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.4))
    p_lbl = lbl_box.text_frame.paragraphs[0]
    p_lbl.text = "COUPA · SR. LEAD, P&T OPERATIONS · 90-DAY ROLLOUT"
    p_lbl.font.name = 'Segoe UI'
    p_lbl.font.size = Pt(11)
    p_lbl.font.bold = True
    p_lbl.font.color.rgb = COLOR_BLUE

    # Line
    line_shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.7), Inches(1.5), Inches(0.06))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = COLOR_BLUE
    line_shape.line.color.rgb = COLOR_BLUE

    # Title
    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.6))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "Turning the CuCP Framework\ninto Adopted Reality"
    p1.font.name = 'Segoe UI'
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_DARK

    # Subtitle
    sub_box = slide1.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(10.0), Inches(0.8))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "A 90-day plan to implement and scale the Customer Collaboration Program — from paper framework to predictable launch engine."
    p_sub.font.name = 'Segoe UI'
    p_sub.font.size = Pt(15)
    p_sub.font.color.rgb = COLOR_MUTED

    # Meta Section
    meta_items = [
        ("Presented by", "Vittobha Vignesh S"),
        ("Role", "Sr. Lead, P&T Operations"),
        ("Programme", "Customer Collaboration (CuCP)"),
        ("Scope", "Days 1–90")
    ]
    for i, (label, val) in enumerate(meta_items):
        box = slide1.shapes.add_textbox(Inches(0.8 + i * 2.8), Inches(5.2), Inches(2.6), Inches(0.8))
        tf = box.text_frame
        tf.word_wrap = True
        
        p_l = tf.paragraphs[0]
        p_l.text = label.upper()
        p_l.font.name = 'Segoe UI'
        p_l.font.size = Pt(10)
        p_l.font.bold = True
        p_l.font.color.rgb = COLOR_MUTED
        
        p_v = tf.add_paragraph()
        p_v.text = val
        p_v.font.name = 'Segoe UI'
        p_v.font.size = Pt(13)
        p_v.font.bold = True
        p_v.font.color.rgb = COLOR_DARK

    # ==========================================
    # SLIDE 2: Change Management & Adoption Strategy
    # ==========================================
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide2)
    add_header(slide2, "Change Management & Adoption Strategy", "Change management starts with understanding the resistance, then flipping the narrative.")

    # Left Card (Problem)
    card_l = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(3.2))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = COLOR_WHITE
    card_l.line.color.rgb = COLOR_RED
    card_l.line.width = Pt(1.5)
    
    tb_l = slide2.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.3), Inches(2.8))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "THE PROBLEM: FRAGMENTED EXECUTION"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED
    
    bullets_l = [
        "Ad-hoc EAP/LA engagement, no standard process",
        "\"Definition drift\" — scope shifts before GA",
        "No exit criteria → risky launch decisions",
        "PMs see CuCP as admin overhead slowing velocity"
    ]
    for b in bullets_l:
        p_b = tf_l.add_paragraph()
        p_b.text = "•  " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = COLOR_DARK
        p_b.space_after = Pt(8)

    # Right Card (Strategy)
    card_r = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(3.2))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = COLOR_WHITE
    card_r.line.color.rgb = COLOR_GREEN
    card_r.line.width = Pt(1.5)
    
    tb_r = slide2.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(2.8))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "STRATEGY: EARN ADOPTION, DON'T MANDATE IT"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN
    
    bullets_r = [
        "Lead with PM pain — not programme goals",
        "Find proto-champions already doing EAP work",
        "Make first experience frictionless — I absorb the admin",
        "Their story converts sceptics, mine doesn't"
    ]
    for b in bullets_r:
        p_b = tf_r.add_paragraph()
        p_b.text = "•  " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = COLOR_DARK
        p_b.space_after = Pt(8)

    # Bottom Reframe Card
    card_b = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.4))
    card_b.fill.solid()
    card_b.fill.fore_color.rgb = COLOR_BLUE_LIGHT
    card_b.line.color.rgb = COLOR_BLUE
    card_b.line.width = Pt(1)
    
    tb_b = slide2.shapes.add_textbox(Inches(1.0), Inches(5.3), Inches(11.3), Inches(1.2))
    tf_b = tb_b.text_frame
    tf_b.word_wrap = True
    p_ref = tf_b.paragraphs[0]
    p_ref.text = "The message that lands: \"CuCP is your evidence kit — not a compliance checklist.\""
    p_ref.font.name = 'Segoe UI'
    p_ref.font.size = Pt(15)
    p_ref.font.bold = True
    p_ref.font.color.rgb = COLOR_BLUE
    
    p_ref2 = tf_b.add_paragraph()
    p_ref2.text = "Structured EAP data protects PMs at GA. It answers \"how do we know this is ready?\" with customer signal, not gut feel."
    p_ref2.font.name = 'Segoe UI'
    p_ref2.font.size = Pt(13)
    p_ref2.font.color.rgb = COLOR_DARK

    # ==========================================
    # SLIDE 3: Capability Classifications: Standard vs. AI
    # ==========================================
    slide3 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide3)
    add_header(slide3, "Capability Classifications: Standard vs. AI", "Pre-GA validation strategies are structured to match the risk profile of each release path.")

    # Left Column: Standard Capabilities
    card_c1 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(3.2))
    card_c1.fill.solid()
    card_c1.fill.fore_color.rgb = COLOR_WHITE
    card_c1.line.color.rgb = COLOR_BLUE
    card_c1.line.width = Pt(1.5)

    tb_c1 = slide3.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(5.3), Inches(2.9))
    tf_c1 = tb_c1.text_frame
    tf_c1.word_wrap = True
    
    p = tf_c1.paragraphs[0]
    p.text = "STANDARD CAPABILITIES"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE
    
    bullets_c1 = [
        "Definition: Deterministic logic and workflows (no AI involved). E.g., dashboards, routing rules.",
        "Focus: Validate operational readiness, scalability, support enablement.",
        "Gating: Closed Beta (EAP) and Limited Availability (LA) require a $0 Order Form.",
        "Exit: LA phase gates readiness with commercial alignment before General Availability."
    ]
    for b in bullets_c1:
        p_b = tf_c1.add_paragraph()
        p_b.text = "•  " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(12.5)
        p_b.font.color.rgb = COLOR_DARK
        p_b.space_before = Pt(6)

    # Right Column: AI Capabilities
    card_c2 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(3.2))
    card_c2.fill.solid()
    card_c2.fill.fore_color.rgb = COLOR_WHITE
    card_c2.line.color.rgb = COLOR_GREEN
    card_c2.line.width = Pt(1.5)

    tb_c2 = slide3.shapes.add_textbox(Inches(7.0), Inches(1.9), Inches(5.3), Inches(2.9))
    tf_c2 = tb_c2.text_frame
    tf_c2.word_wrap = True
    
    p = tf_c2.paragraphs[0]
    p.text = "AI CAPABILITIES (AI-ASSISTED & AGENTIC)"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN
    
    bullets_c2 = [
        "Definition: Powered by ML/LLMs. Probabilistic models. E.g., Copilots, Agent Studio.",
        "Focus: Model tuning, trust/safety guardrails, telemetry integration.",
        "Billing: EAP & Open Beta are free of charge via customer AI Credits.",
        "Gating: Open Beta is enabled via UI toggle (no $0 Order Forms; utilizes Fair Use Caps)."
    ]
    for b in bullets_c2:
        p_b = tf_c2.add_paragraph()
        p_b.text = "•  " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(12.5)
        p_b.font.color.rgb = COLOR_DARK
        p_b.space_before = Pt(6)

    # Bottom Proportional Governance Banner
    card_gov = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.4))
    card_gov.fill.solid()
    card_gov.fill.fore_color.rgb = COLOR_BLUE_LIGHT
    card_gov.line.color.rgb = COLOR_BLUE
    card_gov.line.width = Pt(1)

    tb_gov = slide3.shapes.add_textbox(Inches(1.0), Inches(5.3), Inches(11.3), Inches(1.2))
    tf_gov = tb_gov.text_frame
    tf_gov.word_wrap = True
    
    p = tf_gov.paragraphs[0]
    p.text = "THE PRINCIPLE OF PROPORTIONAL GOVERNANCE"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE
    
    p_g2 = tf_gov.add_paragraph()
    p_g2.text = "The level of release gating must align to feature risk, complexity, and operational impact. Pre-GA (LA & Open Beta) phases are optional and variable. The PM determines if Phase 2 is required based on features."
    p_g2.font.name = 'Segoe UI'
    p_g2.font.size = Pt(12)
    p_g2.font.color.rgb = COLOR_DARK
    p_g2.space_before = Pt(4)


    # ==========================================
    # SLIDE 4: 90-Day Rollout Execution
    # ==========================================
    slide4 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide4)
    add_header(slide4, "90-Day Rollout Execution: Pilot → Learn → Scale", "No big-bang mandate. Build trust through service, then scale on proven evidence.")

    phases = [
        ("Phase 1", "Days 1–30", "Listen & Align", COLOR_GOLD, [
            "1:1 listening sessions (8–10 PMs)",
            "Shadow EAP/LA release cycles",
            "Map informal processes in use",
            "Identify proto-champions",
            "Select pilots: 1 Std + 1 AI release"
        ], "Friction map + pilot cohort confirmed", COLOR_GOLD_LIGHT, COLOR_DARK),
        ("Phase 2", "Days 31–60", "Pilot & Refine", COLOR_BLUE, [
            "Run 2–3 live CuCP pilots",
            "I act as embedded CuCP Concierge",
            "Weekly retros — cut admin overhead",
            "Capture before/after evidence",
            "Activate PM champions"
        ], "Playbook v1.0 + champion stories", COLOR_BLUE_LIGHT, COLOR_BLUE),
        ("Phase 3", "Days 61–90", "Scale & Institutionalise", COLOR_GREEN, [
            "All-hands enablement session",
            "Self-serve KB + templates live",
            "#cucp-support + office hours",
            "Admin tasks automated",
            "Leadership dashboard live"
        ], "All PMs enabled · tracking live", COLOR_GREEN, COLOR_WHITE)
    ]

    for idx, (p_num, p_days, p_goal, p_color, p_bullets, p_out, p_out_bg, p_out_fg) in enumerate(phases):
        # Phase card shape
        p_card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + idx * 4.0), Inches(1.8), Inches(3.7), Inches(4.8))
        p_card.fill.solid()
        p_card.fill.fore_color.rgb = COLOR_WHITE
        p_card.line.color.rgb = p_color
        p_card.line.width = Pt(2)
        
        tb = slide4.shapes.add_textbox(Inches(0.95 + idx * 4.0), Inches(1.9), Inches(3.4), Inches(4.6))
        tf = tb.text_frame
        tf.word_wrap = True
        
        # Num
        p = tf.paragraphs[0]
        p.text = p_num.upper()
        p.font.name = 'Segoe UI'
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = p_color
        
        # Days
        p = tf.add_paragraph()
        p.text = p_days
        p.font.name = 'Segoe UI'
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_DARK
        
        # Goal
        p = tf.add_paragraph()
        p.text = p_goal
        p.font.name = 'Segoe UI'
        p.font.size = Pt(13)
        p.font.italic = True
        p.font.color.rgb = COLOR_MUTED
        p.space_after = Pt(8)
        
        # Bullets
        for b in p_bullets:
            p_b = tf.add_paragraph()
            p_b.text = "•  " + b
            p_b.font.name = 'Segoe UI'
            p_b.font.size = Pt(12)
            p_b.font.color.rgb = COLOR_DARK
            p_b.space_after = Pt(4)
        
        # Output pill shape inside card
        p_pill = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95 + idx * 4.0), Inches(5.7), Inches(3.4), Inches(0.7))
        p_pill.fill.solid()
        p_pill.fill.fore_color.rgb = p_out_bg
        p_pill.line.fill.background()
        
        tb_pill = slide4.shapes.add_textbox(Inches(0.95 + idx * 4.0), Inches(5.7), Inches(3.4), Inches(0.7))
        tf_pill = tb_pill.text_frame
        tf_pill.word_wrap = True
        p_p = tf_pill.paragraphs[0]
        p_p.text = p_out
        p_p.alignment = PP_ALIGN.CENTER
        p_p.font.name = 'Segoe UI'
        p_p.font.size = Pt(11)
        p_p.font.bold = True
        p_p.font.color.rgb = p_out_fg

    # ==========================================
    # SLIDE 5: Enablement & Support Model
    # ==========================================
    slide5 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide5)
    add_header(slide5, "Enablement & Support Model", "Four pillars that make CuCP stupidly easy to run — before it becomes expected.")

    pillars = [
        ("🎯", "CuCP Concierge (Days 1–60)", "I act as the operational support for pilot PMs. They learn the process; I handle templates, customer invites, and exit criteria tracking. PM overhead: near zero.", "High-touch during pilot"),
        ("📋", "Playbook + Cheat Sheets", "Two 1-page cheat sheets: Standard capability path and AI capability path. Separate because their pre-GA goals, cohort sizes, and exit criteria differ meaningfully.", "Self-serve from Day 60"),
        ("💬", "#cucp-support + Office Hours", "Live Slack channel for real-time questions. Weekly 30-min open office hours for the first 90 days. PMs get answers in minutes, not days — reduces perceived friction.", "Async + sync support"),
        ("🏆", "Champion-Led Rollout", "Pilot PMs present at the all-hands enablement session — not me. Peer stories are far more credible than top-down messaging for converting sceptics.", "Peer credibility > authority")
    ]

    for idx, (icon, title, desc, tag) in enumerate(pillars):
        col = idx % 2
        row = idx // 2
        
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + col * 6.0), Inches(1.8 + row * 2.5), Inches(5.7), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = COLOR_BLUE_LIGHT
        card.line.width = Pt(1.5)
        
        tb = slide5.shapes.add_textbox(Inches(1.0 + col * 6.0), Inches(1.9 + row * 2.5), Inches(5.3), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"{icon}  {title}"
        p.font.name = 'Segoe UI'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_DARK
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = 'Segoe UI'
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = COLOR_MUTED
        p_d.space_before = Pt(6)
        
        # Tag shape
        tag_sh = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0 + col * 6.0), Inches(3.4 + row * 2.5), Inches(2.5), Inches(0.35))
        tag_sh.fill.solid()
        tag_sh.fill.fore_color.rgb = COLOR_BLUE_LIGHT
        tag_sh.line.fill.background()
        
        tb_tag = slide5.shapes.add_textbox(Inches(1.0 + col * 6.0), Inches(3.4 + row * 2.5), Inches(2.5), Inches(0.35))
        tf_t = tb_tag.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = tag
        p_t.alignment = PP_ALIGN.CENTER
        p_t.font.name = 'Segoe UI'
        p_t.font.size = Pt(10)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_BLUE

    # ==========================================
    # SLIDE 6: Measuring Initial Success
    # ==========================================
    slide6 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide6)
    add_header(slide6, "Measuring Initial Success", "Six metrics prove the rollout is working. One leading indicator matters most in the first 90 days.")

    # Table coordinates
    left = Inches(0.8)
    top = Inches(1.8)
    width = Inches(11.7)
    height = Inches(4.5)
    
    table_shape = slide6.shapes.add_table(7, 5, left, top, width, height)
    table = table_shape.table
    
    # Column widths
    table.columns[0].width = Inches(3.8) # KPI
    table.columns[1].width = Inches(1.6) # Baseline
    table.columns[2].width = Inches(1.8) # Day 90 Target
    table.columns[3].width = Inches(2.9) # How Measured
    table.columns[4].width = Inches(1.6) # Type
    
    # Headers
    headers = ["KPI", "Baseline", "Day 90 Target", "How Measured", "Type"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(244, 245, 247)
        p = cell.text_frame.paragraphs[0]
        p.text = h.upper()
        p.alignment = PP_ALIGN.LEFT
        p.font.name = 'Segoe UI'
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_MUTED

    rows_data = [
        ("⭐ PM Satisfaction Score — \"Is CuCP helping?\"", "—", "≥ 4.0 / 5.0", "Post-pilot survey", "LEADING"),
        ("% eligible releases using CuCP", "0% (ad-hoc)", "≥ 50%", "Release tracker", "LAGGING"),
        ("EAP → GA cycle time", "Measure Days 1–30", "No regression; –10%", "Release logs", "LAGGING"),
        ("Definition drift incidents", "Unmeasured", "Tracked; 0 in pilots", "EAP/LA retros", "LEADING"),
        ("Customer participation rate", "—", "≥ 70% invited engage", "CuCP tracker", "LAGGING"),
        ("PM admin time per release cycle", "—", "< 2 hours", "PM self-report", "LEADING")
    ]

    for row_idx, row in enumerate(rows_data):
        is_star = row[0].startswith("⭐")
        for col_idx, text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            if is_star:
                cell.fill.fore_color.rgb = COLOR_GOLD_LIGHT
            else:
                cell.fill.fore_color.rgb = COLOR_WHITE
            
            p = cell.text_frame.paragraphs[0]
            p.text = text
            p.font.name = 'Segoe UI'
            p.font.size = Pt(11.5)
            p.font.color.rgb = COLOR_DARK
            if is_star:
                p.font.bold = True
            
            # Highlight type tag
            if col_idx == 4:
                p.font.bold = True
                if text == "LEADING":
                    p.font.color.rgb = COLOR_GREEN
                else:
                    p.font.color.rgb = COLOR_BLUE

    # Footnote
    fn_box = slide6.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.4))
    p_fn = fn_box.text_frame.paragraphs[0]
    p_fn.text = "⭐ Key takeaway: If PMs feel CuCP helps — scaling is pull, not push. Fix the playbook before mandating."
    p_fn.font.name = 'Segoe UI'
    p_fn.font.size = Pt(12)
    p_fn.font.bold = True
    p_fn.font.color.rgb = COLOR_GOLD

    # ==========================================
    # SLIDE 7: Day 90 Vision & Ask
    # ==========================================
    slide7 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide7)
    add_header(slide7, "What Day 90 Looks Like", "Three outcomes that prove the CuCP rollout succeeded — and what I need from leadership to get there.")

    # Left box: Evidence
    card_v1 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(2.9))
    card_v1.fill.solid()
    card_v1.fill.fore_color.rgb = COLOR_WHITE
    card_v1.line.color.rgb = COLOR_BLUE_LIGHT
    card_v1.line.width = Pt(1.5)
    
    tb_v1 = slide7.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(5.3), Inches(2.7))
    tf_v1 = tb_v1.text_frame
    tf_v1.word_wrap = True
    p = tf_v1.paragraphs[0]
    p.text = "EVIDENCE OF SUCCESS"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE
    
    bullets_v1 = [
        "≥ 50% eligible releases running CuCP",
        "PM satisfaction score ≥ 4.0 / 5.0",
        "2–3 champion PMs advocating internally",
        "Self-serve KB live — zero dependency on me"
    ]
    for b in bullets_v1:
        p_b = tf_v1.add_paragraph()
        p_b.text = "✓  " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = COLOR_DARK
        p_b.space_after = Pt(4)

    # Right box: Governance
    card_v2 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.9))
    card_v2.fill.solid()
    card_v2.fill.fore_color.rgb = COLOR_WHITE
    card_v2.line.color.rgb = COLOR_BLUE_LIGHT
    card_v2.line.width = Pt(1.5)
    
    tb_v2 = slide7.shapes.add_textbox(Inches(7.0), Inches(1.9), Inches(5.3), Inches(2.7))
    tf_v2 = tb_v2.text_frame
    tf_v2.word_wrap = True
    p = tf_v2.paragraphs[0]
    p.text = "PROPORTIONAL GOVERNANCE PRINCIPLE"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE
    
    bullets_v2 = [
        "Days 1–30: Observe only, map release landscape",
        "Days 31–60: Voluntary pilot, embedding Concierge",
        "Days 61–90: Clear expectation for eligible cohorts",
        "Day 90+: Data-backed mandate, confirmed by feedback"
    ]
    for b in bullets_v2:
        p_b = tf_v2.add_paragraph()
        p_b.text = "✓  " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = COLOR_DARK
        p_b.space_after = Pt(4)

    # Bottom Ask Box
    card_ask = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.9), Inches(11.7), Inches(1.8))
    card_ask.fill.solid()
    card_ask.fill.fore_color.rgb = COLOR_BLUE_LIGHT
    card_ask.line.color.rgb = COLOR_BLUE
    card_ask.line.width = Pt(1)
    
    tb_ask = slide7.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(11.3), Inches(1.6))
    tf_ask = tb_ask.text_frame
    tf_ask.word_wrap = True
    p = tf_ask.paragraphs[0]
    p.text = "MY ASK OF P&T LEADERSHIP"
    p.font.name = 'Segoe UI'
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE
    
    asks = [
        "🛡️ Air cover during pilot — no forced mandate",
        "🎙️ Exec sponsor for all-hands session",
        "🔗 Access to 2–3 willing pilot PMs in Week 1",
        "📊 Release log access for baseline measurement"
    ]
    
    # We can layout the asks horizontally
    for i, ask in enumerate(asks):
        ask_sh = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0 + i * 2.8), Inches(5.5), Inches(2.7), Inches(0.9))
        ask_sh.fill.solid()
        ask_sh.fill.fore_color.rgb = COLOR_WHITE
        ask_sh.line.color.rgb = COLOR_BLUE_LIGHT
        ask_sh.line.width = Pt(1)
        
        tb_a = slide7.shapes.add_textbox(Inches(1.05 + i * 2.8), Inches(5.5), Inches(2.6), Inches(0.9))
        tf_a = tb_a.text_frame
        tf_a.word_wrap = True
        p_a = tf_a.paragraphs[0]
        p_a.text = ask
        p_a.font.name = 'Segoe UI'
        p_a.font.size = Pt(11)
        p_a.font.bold = True
        p_a.font.color.rgb = COLOR_DARK

    # Save presentation
    prs.save(output_path)
    print(f"Presentation successfully created at: {output_path}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        out = "docs/cucp/cucp_presentation.pptx"
    else:
        out = sys.argv[1]
    
    # Ensure dir exists
    os.makedirs(os.path.dirname(os.path.abspath(out)), exist_ok=True)
    create_presentation(out)
